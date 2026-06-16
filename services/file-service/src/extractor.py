"""
In-memory text extraction for uploaded documents (PDF / DOCX / PPTX).

Nothing is ever written to disk: the raw bytes are parsed straight from memory
(io.BytesIO) and are discarded by the caller as soon as extraction returns.
"""
import io
import logging
import re
from pathlib import Path

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx"}

class UnsupportedDocumentError(ValueError):
    """Raised when a file extension is not one of the supported document types."""

class DocumentExtractionError(RuntimeError):
    """Raised when a supported file cannot be parsed."""

def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:  # noqa: BLE001 — skip an individual unreadable page
            continue
    return "\n".join(parts)


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text for c in row.cells if c.text]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs)
                if text.strip():
                    parts.append(text)
    return "\n".join(parts)


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
}


def _normalise(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_text(filename: str, data: bytes, *, max_chars: int | None = None) -> str:
    """
    Extract plain text from an in-memory document.

    Raises UnsupportedDocumentError for unknown extensions and
    DocumentExtractionError when a supported file cannot be parsed.
    """
    ext = Path(filename or "").suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise UnsupportedDocumentError(f"Unsupported file type: '{ext or filename}'")

    try:
        text = extractor(data)
    except Exception as exc:  # noqa: BLE001 — surface a single, typed error
        raise DocumentExtractionError(f"Failed to parse {ext} document: {exc}") from exc

    text = _normalise(text)
    if max_chars is not None and len(text) > max_chars:
        text = text[:max_chars]
    return text


def suggest_topic(text: str, filename: str = "") -> str:
    """Derive a short, human-readable quiz title from the file name or first line."""
    stem = Path(filename or "").stem.replace("_", " ").replace("-", " ").strip()
    if stem:
        return stem[:200]
    first_line = next((ln.strip() for ln in text.splitlines() if ln.strip()), "")
    return first_line[:200] or "Documento caricato"
